from functools import wraps
from flask import jsonify
from flask_jwt_extended import jwt_required, get_jwt_identity
from app.models import User
import PyPDF2
from io import BytesIO
from docx import Document # Import Document for docx files
from pptx import Presentation # Import Presentation for pptx files

def role_required(role_name):
    def wrapper(fn):
        @wraps(fn) # Preserve original function's metadata
        @jwt_required()
        def decorator(*args, **kwargs):
            current_user_id = get_jwt_identity()
            user = User.query.get(current_user_id)
            if user and user.role.name == role_name:
                return fn(*args, **kwargs)
            return jsonify({"msg": f"'{role_name}' role required"}), 403
        return decorator
    return wrapper

def extract_text_from_file(file_content_bytes, file_type):
    """
    Extracts text content from a file (PDF or plain text).
    :param file_content_bytes: Bytes content of the file.
    :param file_type: MIME type of the file (e.g., 'application/pdf', 'text/plain').
    :return: Extracted text as a string.
    """
    if file_type == 'application/pdf':
        try:
            pdf_file = BytesIO(file_content_bytes)
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            text = ""
            for page_num in range(len(pdf_reader.pages)):
                text += pdf_reader.pages[page_num].extract_text() or ""
            return text
        except Exception as e:
            raise ValueError(f"Error extracting text from PDF: {e}")
    elif file_type.startswith('text/'):
        try:
            return file_content_bytes.decode('utf-8')
        except UnicodeDecodeError:
            return file_content_bytes.decode('latin-1') # Fallback for other encodings
    elif file_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document': # .docx files
        try:
            doc = Document(BytesIO(file_content_bytes))
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text
        except Exception as e:
            raise ValueError(f"Error extracting text from DOCX: {e}")
    elif file_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation': # .pptx files
        try:
            ppt = Presentation(BytesIO(file_content_bytes))
            text = ""
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except Exception as e:
            raise ValueError(f"Error extracting text from PPTX: {e}")
    else:
        raise ValueError(f"Unsupported file type for text extraction: {file_type}")
