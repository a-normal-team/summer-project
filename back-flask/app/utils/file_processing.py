"""
文件处理相关的工具函数
包括文本提取、文件解析等功能
"""

import PyPDF2
from io import BytesIO
from docx import Document  # 用于处理docx文件
from pptx import Presentation  # 用于处理pptx文件

def extract_text_from_file(file_content_bytes, file_type):
    """
    从文件（PDF、文本文件等）中提取文本内容
    :param file_content_bytes: 文件的字节内容
    :param file_type: 文件的MIME类型（例如：'application/pdf', 'text/plain'）
    :return: 提取的文本（字符串）
    """
    if file_type == 'application/pdf':
        try:
            pdf_file = BytesIO(file_content_bytes)
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            text = ""
            for page_num in range(len(pdf_reader.pages)):
                text += pdf_reader.pages[page_num].extract_text() or ""
            return text
        except Exception as e:
            raise ValueError(f"Error extracting text from PDF: {e}")
    elif file_type.startswith('text/'):
        # 对于文本文件，直接解码
        try:
            return file_content_bytes.decode('utf-8')
        except UnicodeDecodeError:
            try:
                # 尝试其他编码
                return file_content_bytes.decode('latin-1')
            except Exception as e:
                raise ValueError(f"Error decoding text file: {e}")
    elif file_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
        # 处理docx文件
        try:
            docx_file = BytesIO(file_content_bytes)
            doc = Document(docx_file)
            return "\n".join(paragraph.text for paragraph in doc.paragraphs)
        except Exception as e:
            raise ValueError(f"Error extracting text from DOCX: {e}")
    elif file_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
        # 处理pptx文件
        try:
            pptx_file = BytesIO(file_content_bytes)
            prs = Presentation(pptx_file)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except Exception as e:
            raise ValueError(f"Error extracting text from PPTX: {e}")
    else:
        raise ValueError(f"Unsupported file type: {file_type}")
